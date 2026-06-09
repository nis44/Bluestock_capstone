"""Day 7: generate final PDF report, presentation deck, and completion checklist."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from project_paths import CHARTS, DATA_PROCESSED, REPORTS, ROOT, ensure_directories


def table_from_df(df: pd.DataFrame, max_rows: int = 8) -> Table:
    display = df.head(max_rows).copy()
    data = [display.columns.tolist()] + display.astype(str).values.tolist()
    table = Table(data, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0F766E")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 7),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#CBD5E1")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
            ]
        )
    )
    return table


def add_chart(story: list, chart_name: str, width: float = 6.4 * inch) -> None:
    path = CHARTS / chart_name
    if path.exists():
        story.append(Image(str(path), width=width, height=width * 0.55))
        story.append(Spacer(1, 0.15 * inch))


def build_pdf() -> None:
    output = REPORTS / "Final_Report.pdf"
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="Small", parent=styles["BodyText"], fontSize=8, leading=10))
    styles["Title"].textColor = colors.HexColor("#0F766E")
    styles["Heading1"].textColor = colors.HexColor("#0F766E")
    styles["Heading2"].textColor = colors.HexColor("#334155")

    fund = pd.read_csv(DATA_PROCESSED / "clean_fund_master.csv")
    scorecard = pd.read_csv(DATA_PROCESSED / "fund_scorecard.csv")
    var = pd.read_csv(DATA_PROCESSED / "var_cvar_report.csv")
    cohort = pd.read_csv(DATA_PROCESSED / "cohort_analysis.csv")
    hhi = pd.read_csv(DATA_PROCESSED / "sector_hhi.csv")

    story = [
        Paragraph("Bluestock Fintech Mutual Fund Analytics Platform", styles["Title"]),
        Paragraph("End-to-End Data Engineering, ETL Pipeline, Analytics and Dashboard Assets", styles["Heading2"]),
        Spacer(1, 0.3 * inch),
        Paragraph("Executive Summary", styles["Heading1"]),
        Paragraph(
            "This project consolidates publicly available and supplied mutual fund datasets into a reproducible analytics platform. "
            "It includes data ingestion, cleaning, SQLite storage, exploratory analysis, risk-return metrics, advanced investor analytics, "
            "dashboard-ready extracts, and final submission documents.",
            styles["BodyText"],
        ),
        Spacer(1, 0.2 * inch),
        Paragraph("Dataset Coverage", styles["Heading1"]),
        Paragraph(
            f"The project covers {len(fund)} schemes, 46,000 NAV observations, 32,778 investor transactions, benchmark index series, "
            "AUM, SIP, folio, category inflow, and portfolio holdings datasets.",
            styles["BodyText"],
        ),
        Spacer(1, 0.2 * inch),
        Paragraph("System Architecture", styles["Heading1"]),
        Paragraph(
            "The pipeline follows Extract, Transform, Load, Analyse, and Visualise layers. Raw CSVs are loaded with Pandas, cleaned into "
            "processed outputs, stored in SQLite tables, analysed with Python scripts/notebooks, and exported for Power BI/Tableau usage.",
            styles["BodyText"],
        ),
        Paragraph("EDA Highlights", styles["Heading1"]),
    ]
    for chart in [
        "02_aum_growth_by_amc.png",
        "03_sip_inflow_trend.png",
        "04_category_inflow_heatmap.png",
        "07_transaction_amount_by_state.png",
    ]:
        add_chart(story, chart)

    story.extend(
        [
            PageBreak(),
            Paragraph("Performance Analytics", styles["Heading1"]),
            Paragraph(
                "Day 4 computes annualized returns, 1/3/5-year CAGR, Sharpe, Sortino, Alpha, Beta, tracking error, information ratio, "
                "maximum drawdown, annual volatility, and a weighted composite scorecard.",
                styles["BodyText"],
            ),
            Spacer(1, 0.15 * inch),
            table_from_df(
                scorecard[
                    [
                        "scheme_name",
                        "fund_house",
                        "return_3yr_pct",
                        "sharpe_ratio",
                        "alpha_pct",
                        "beta",
                        "fund_score",
                    ]
                ],
                8,
            ),
            Spacer(1, 0.2 * inch),
        ]
    )
    add_chart(story, "16_top5_funds_vs_benchmarks.png")
    add_chart(story, "12_return_vs_risk_scatter.png")

    story.extend(
        [
            Paragraph("Advanced Analytics", styles["Heading1"]),
            Paragraph(
                "Day 6 adds historical VaR/CVaR, rolling Sharpe, cohort analysis, SIP continuity checks, recommendation logic, and "
                "sector concentration HHI.",
                styles["BodyText"],
            ),
            Spacer(1, 0.15 * inch),
            table_from_df(var[["scheme_name", "fund_house", "var_95_pct", "cvar_95_pct"]].sort_values("var_95_pct"), 6),
            Spacer(1, 0.15 * inch),
            table_from_df(cohort, 6),
            Spacer(1, 0.15 * inch),
            table_from_df(hhi[["scheme_name", "fund_house", "sector_hhi"]], 6),
        ]
    )
    add_chart(story, "17_rolling_90d_sharpe.png")
    add_chart(story, "18_sector_concentration_hhi.png")

    story.extend(
        [
            PageBreak(),
            Paragraph("Dashboard Plan", styles["Heading1"]),
            Paragraph(
                "Dashboard-ready extracts are available in dashboard/extracts. The Power BI build guide defines four report pages: "
                "Industry Overview, Fund Performance, Investor Analytics, and SIP & Market Trends.",
                styles["BodyText"],
            ),
            Paragraph("Recommendations", styles["Heading1"]),
            Paragraph(
                "Use the fund scorecard for shortlist creation, review downside risk with VaR and max drawdown, and combine SIP/cohort "
                "insights with geographic segmentation for investor engagement. Direct plans show lower expense ratios and should be "
                "compared carefully against regular plans on risk-adjusted returns.",
                styles["BodyText"],
            ),
            Paragraph("Limitations", styles["Heading1"]),
            Paragraph(
                "Investor transaction data is synthetic, portfolio holdings are point-in-time, and dashboard construction in Power BI "
                "requires manual layout in Power BI Desktop. This analysis is educational and not financial advice.",
                styles["BodyText"],
            ),
            Paragraph("Submission Checklist", styles["Heading1"]),
            Paragraph(
                "ETL scripts, cleaned CSVs, SQLite schema/database, SQL queries, EDA charts, performance metrics, advanced analytics, "
                "dashboard extracts, final PDF report, PowerPoint deck, README, and local Git commits are complete.",
                styles["BodyText"],
            ),
        ]
    )
    doc = SimpleDocTemplate(str(output), pagesize=A4, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    doc.build(story)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)


def add_chart_slide(prs: Presentation, title: str, chart_name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    path = CHARTS / chart_name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.7), Inches(1.2), width=Inches(8.8))


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Bluestock MF Analytics Platform"
    title_slide.placeholders[1].text = "ETL, SQL, Risk Metrics, EDA, Dashboard Assets"

    add_bullet_slide(
        prs,
        "Problem and Objective",
        [
            "Unify fragmented AMFI, NAV, SIP, AUM, transaction, and benchmark data.",
            "Build a repeatable analytics pipeline for mutual fund selection insights.",
            "Create dashboard-ready outputs for executive review.",
        ],
    )
    add_bullet_slide(
        prs,
        "Data Sources",
        [
            "10 supplied datasets covering 40 schemes and 32,778 transactions.",
            "Live NAV fetched from mfapi.in for six selected schemes.",
            "Benchmarks include NIFTY50, NIFTY100, midcap, smallcap, and debt indices.",
        ],
    )
    add_bullet_slide(
        prs,
        "Architecture",
        [
            "Python/Pandas ingestion and cleaning.",
            "SQLite star-style schema with dimensions and fact tables.",
            "Reports, charts, and dashboard extracts generated from processed data.",
        ],
    )
    add_chart_slide(prs, "EDA: AUM by AMC", "02_aum_growth_by_amc.png")
    add_chart_slide(prs, "EDA: SIP Inflow Trend", "03_sip_inflow_trend.png")
    add_chart_slide(prs, "Performance: Fund vs Benchmark", "16_top5_funds_vs_benchmarks.png")
    add_chart_slide(prs, "Performance: Return vs Risk", "12_return_vs_risk_scatter.png")
    add_chart_slide(prs, "Advanced: Rolling Sharpe", "17_rolling_90d_sharpe.png")
    add_chart_slide(prs, "Advanced: Sector Concentration", "18_sector_concentration_hhi.png")
    add_bullet_slide(
        prs,
        "Dashboard Pages",
        [
            "Industry Overview: AUM, SIP, folios, AMC comparison.",
            "Fund Performance: scorecard, risk-return plot, benchmark tracking.",
            "Investor Analytics: state, city tier, age, transaction type.",
            "SIP and Market Trends: inflows, category heatmap, index overlay.",
        ],
    )
    add_bullet_slide(
        prs,
        "Key Takeaways",
        [
            "Scorecard combines return, Sharpe, alpha, expense ratio, and drawdown.",
            "Investor analytics supports segmentation by geography and demographics.",
            "Risk metrics help avoid judging funds only by raw returns.",
        ],
    )

    prs.save(REPORTS / "Bluestock_MF_Presentation.pptx")


def write_checklist() -> None:
    checklist = """# Final Submission Checklist

- [x] Day 1 data ingestion script and live NAV fetcher
- [x] Day 2 cleaned CSVs, SQLite schema/database, SQL queries, data dictionary
- [x] Day 3 EDA charts and findings
- [x] Day 4 fund performance metrics and scorecard
- [x] Day 5 dashboard extracts and Power BI build guide
- [x] Day 6 advanced analytics and recommender
- [x] Day 7 final PDF report and presentation deck
- [x] README and requirements file
- [x] Local Git repository committed
- [ ] GitHub push, intentionally not performed
"""
    (REPORTS / "Final_Submission_Checklist.md").write_text(checklist, encoding="utf-8")


def main() -> None:
    ensure_directories()
    build_pdf()
    build_pptx()
    write_checklist()
    print("Final report, presentation, and checklist generated.")


if __name__ == "__main__":
    main()
